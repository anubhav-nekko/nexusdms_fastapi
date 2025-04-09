import os
import io
import json
import pickle
import tempfile
from typing import List, Dict, Any, Tuple

import faiss # type: ignore
import fitz  # type: ignore # PyMuPDF
import numpy as np # type: ignore
import pandas as pd # type: ignore
import requests # type: ignore
from fastapi import FastAPI, File, UploadFile, HTTPException # type: ignore
from fastapi.responses import JSONResponse # type: ignore
from pydantic import BaseModel # type: ignore
from sentence_transformers import SentenceTransformer # type: ignore
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions # type: ignore
from docx import Document # type: ignore
from pptx import Presentation # type: ignore
from PIL import Image # type: ignore

from fastapi.middleware.cors import CORSMiddleware # type: ignore

###############################################################################
# 1) LOAD SECRETS AND CONFIG
###############################################################################

SECRETS_PATH = "../secrets.json"

def load_secrets(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)

SECRETS = load_secrets(SECRETS_PATH)

connection_string = SECRETS["connection_string"]  # For Azure Blob
azure_account_name = SECRETS["azure_account_name"]
azure_account_key = SECRETS["azure_account_key"]
container_name = SECRETS["container_name"]

GPT_ENDPOINT = SECRETS["GPT_ENDPOINT"]
GPT_API = SECRETS["GPT_API"]

FAISS_INDEX_PATH = SECRETS["FAISS_INDEX_PATH"]       # e.g., "faiss_index.index"
METADATA_STORE_PATH = SECRETS["METADATA_STORE_PATH"] # e.g., "metadata.pkl"

# Initialize a sentence-transformers model
mpnet_model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")

# Global system prompt (tweak as needed)
SYSTEM_MESSAGE = """
You are an advanced legal data analyst specializing in legal document analysis. 
Provide an in-depth, accurate answer using only the given context. 
If information is missing or unclear, say so explicitly.
"""

###############################################################################
# 2) FASTAPI APP
###############################################################################
app = FastAPI(title="Azure Document Query Assistant", version="1.0.0")

# Allow all origins (not recommended for production)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Or specify a list like ["https://yourfrontend.com"]
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

###############################################################################
# 3) FAISS + METADATA INIT
###############################################################################
DIMENSION = 768  # MPNet dimension
faiss_index = faiss.IndexFlatL2(DIMENSION)  # A simple L2 index
metadata_store: List[dict] = []  # Will hold info like { filename, page, text }

###############################################################################
# 4) AZURE BLOB STORAGE UTILS
###############################################################################
def file_exists_in_blob(file_name: str) -> bool:
    """Check if the file_name exists in the Azure Blob container."""
    blob_service_client = BlobServiceClient.from_connection_string(connection_string)
    container_client = blob_service_client.get_container_client(container_name)
    blob_list = [blob.name for blob in container_client.list_blobs()]
    return file_name in blob_list

def upload_to_blob_storage(local_file_path: str, container: str, blob_name: str):
    """Upload a local file to Azure Blob Storage."""
    try:
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)
        blob_client = blob_service_client.get_blob_client(container=container, blob=blob_name)
        with open(local_file_path, "rb") as data:
            blob_client.upload_blob(data, overwrite=True)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to upload to blob: {e}")

def download_from_blob_storage(container: str, blob_name: str, local_file_path: str) -> bool:
    """Download a file from Azure Blob Storage. Returns True if successful, else False if not found."""
    try:
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)
        blob_client = blob_service_client.get_blob_client(container=container, blob=blob_name)
        with open(local_file_path, "wb") as file:
            file.write(blob_client.download_blob().readall())
        return True
    except Exception as e:
        # If 404, return False. Otherwise raise the error.
        if "BlobNotFound" in str(e):
            return False
        raise HTTPException(status_code=500, detail=f"Failed to download from blob: {e}")

###############################################################################
# 5) FAISS INDEX + METADATA LOAD/SAVE
###############################################################################
def save_index_and_metadata():
    """Save the FAISS index + metadata locally, then upload to Azure blob."""
    # Save index locally
    faiss.write_index(faiss_index, FAISS_INDEX_PATH)
    # Save metadata locally
    with open(METADATA_STORE_PATH, "wb") as f:
        pickle.dump(metadata_store, f)

    # Upload both to Azure
    upload_to_blob_storage(FAISS_INDEX_PATH, container_name, os.path.basename(FAISS_INDEX_PATH))
    upload_to_blob_storage(METADATA_STORE_PATH, container_name, os.path.basename(METADATA_STORE_PATH))

def load_index_and_metadata():
    """Attempt to load the FAISS index + metadata from Azure. If not found, start fresh."""
    global faiss_index, metadata_store

    index_blob_name = os.path.basename(FAISS_INDEX_PATH)
    meta_blob_name = os.path.basename(METADATA_STORE_PATH)

    index_ok = download_from_blob_storage(container_name, index_blob_name, FAISS_INDEX_PATH)
    meta_ok = download_from_blob_storage(container_name, meta_blob_name, METADATA_STORE_PATH)

    if index_ok and os.path.exists(FAISS_INDEX_PATH):
        try:
            faiss_index = faiss.read_index(FAISS_INDEX_PATH)
        except Exception as e:
            print("Error reading index from disk:", e)
            faiss_index = faiss.IndexFlatL2(DIMENSION)
    else:
        faiss_index = faiss.IndexFlatL2(DIMENSION)

    if meta_ok and os.path.exists(METADATA_STORE_PATH):
        try:
            with open(METADATA_STORE_PATH, "rb") as f:
                loaded = pickle.load(f)
                metadata_store.clear()
                metadata_store.extend(loaded)
        except Exception as e:
            print("Error reading metadata from disk:", e)
            metadata_store.clear()
    else:
        metadata_store.clear()

###############################################################################
# 6) STARTUP: LOAD INDEX & METADATA
###############################################################################
@app.on_event("startup")
def on_start():
    load_index_and_metadata()
    print("Index & metadata loaded at startup.")

###############################################################################
# 7) UTILITY FUNCTIONS FOR DOCUMENT PROCESSING
###############################################################################
def generate_embeddings(text: str) -> np.ndarray:
    """Create and return an mpnet embedding (float32)."""
    emb = mpnet_model.encode(text, normalize_embeddings=True)
    return emb.astype(np.float32)

def pdf_to_text_pages(pdf_path: str) -> List[Tuple[int, str]]:
    """
    Convert each page of a PDF to text using PyMuPDF. 
    This example does not do OCR (it extracts text if it's digitally embedded). 
    For OCR, you'd plug in an external service here.
    """
    results = []
    doc = fitz.open(pdf_path)
    for i, page in enumerate(doc, start=1):
        text = page.get_text()
        results.append((i, text))
    return results

def docx_to_text_chunks(docx_path: str) -> List[Tuple[int, str]]:
    doc = Document(docx_path)
    full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    chunk_size = 1000
    chunks = []
    i = 0
    start = 0
    while start < len(full_text):
        end = min(start + chunk_size, len(full_text))
        chunk_text = full_text[start:end]
        i += 1
        chunks.append((i, chunk_text))
        start = end
    return chunks

def pptx_to_text_chunks(pptx_path: str) -> List[Tuple[int, str]]:
    prs = Presentation(pptx_path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
        slides.append((idx, "\n".join(lines)))
    return slides

def xlsx_to_text_chunks(xlsx_path: str) -> List[Tuple[int, str]]:
    """
    Simple approach: if extension is .csv, read via pd.read_csv. 
    Otherwise read first sheet. 
    Then chunk each 50 rows as one text block.
    """
    ext = os.path.splitext(xlsx_path)[1].lower()
    if ext == ".csv":
        df = pd.read_csv(xlsx_path)
    else:
        df_excel = pd.ExcelFile(xlsx_path)
        first_sheet = df_excel.sheet_names[0]
        df = pd.read_excel(df_excel, sheet_name=first_sheet)

    chunk_size = 50
    results = []
    idx = 0
    for start in range(0, len(df), chunk_size):
        end = start + chunk_size
        chunk_df = df.iloc[start:end]
        text_str = chunk_df.to_string(index=False)
        idx += 1
        results.append((idx, text_str))
    return results

def image_to_text(img_path: str) -> str:
    """
    Stub for images. If you have actual OCR, you'd do that here.
    This sample just returns a placeholder.
    """
    return "[IMAGE OCR NOT IMPLEMENTED]"  # or call an OCR service

###############################################################################
# 8) UPLOAD DOCUMENT ENDPOINT
###############################################################################

class UploadResponse(BaseModel):
    message: str
    filename: str

@app.post("/upload_document", response_model=UploadResponse)
def upload_document(file: UploadFile = File(...)):
    """
    Accepts a file, saves to temp, checks if it exists in Azure blob, 
    if not, processes to extract text, embed in FAISS, updates metadata, 
    and re-uploads the index + metadata.
    """
    try:
        filename = file.filename
        ext = os.path.splitext(filename)[1].lower()
        temp_path = os.path.join(tempfile.gettempdir(), filename)

        # Save uploaded file to temp
        with open(temp_path, "wb") as f:
            f.write(file.file.read())

        # Check if the file is already in blob
        if file_exists_in_blob(filename):
            return UploadResponse(
                message=f"File `{filename}` already exists in Azure. Skipping re-upload.",
                filename=filename
            )
        else:
            # Upload the raw file to Azure
            upload_to_blob_storage(temp_path, container_name, filename)

        # Now parse the file to extract text, embed, store in FAISS
        new_records = []
        if ext == ".pdf":
            pages = pdf_to_text_pages(temp_path)
            for (pg, txt) in pages:
                if txt.strip():
                    emb = generate_embeddings(txt)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({"filename": filename, "page": pg, "text": txt})
        elif ext in [".doc", ".docx"]:
            chunks = docx_to_text_chunks(temp_path)
            for (pg, txt) in chunks:
                if txt.strip():
                    emb = generate_embeddings(txt)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({"filename": filename, "page": pg, "text": txt})
        elif ext == ".pptx":
            slides = pptx_to_text_chunks(temp_path)
            for (pg, txt) in slides:
                if txt.strip():
                    emb = generate_embeddings(txt)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({"filename": filename, "page": pg, "text": txt})
        elif ext in [".xlsx", ".csv"]:
            table_chunks = xlsx_to_text_chunks(temp_path)
            for (pg, txt) in table_chunks:
                if txt.strip():
                    emb = generate_embeddings(txt)
                    faiss_index.add(emb.reshape(1, -1))
                    new_records.append({"filename": filename, "page": pg, "text": txt})
        elif ext in [".jpg", ".jpeg", ".png"]:
            # If you do real OCR, do it here
            txt = image_to_text(temp_path)
            if txt.strip():
                emb = generate_embeddings(txt)
                faiss_index.add(emb.reshape(1, -1))
                new_records.append({"filename": filename, "page": 1, "text": txt})
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

        # Extend the global metadata
        metadata_store.extend(new_records)

        # Save & upload updated index + metadata
        save_index_and_metadata()

        # Cleanup local temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)

        return UploadResponse(
            message=f"Successfully uploaded and processed `{filename}`.",
            filename=filename
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

###############################################################################
# 9) QUERY DOCUMENTS (BASIC PAGE RANGE FILTER)
###############################################################################

class QueryRequest(BaseModel):
    selected_files: List[str]
    selected_page_ranges: Dict[str, Tuple[int, int]]
    prompt: str
    top_k: int = 5

class QueryResponse(BaseModel):
    answer: str
    top_k_metadata: List[Any]

def call_gpt_api(system_msg: str, user_msg: str) -> str:
    """
    Example of an Azure/OpenAI endpoint call. Adjust to your usage:
    - GPT_ENDPOINT is your Azure OpenAI or custom endpoint
    - GPT_API is your key
    """
    url = GPT_ENDPOINT
    headers = {
        "Content-Type": "application/json",
        "api-key": GPT_API
    }
    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": user_msg}
    ]
    payload = {
        "messages": messages,
        "temperature": 0.7,
        "max_tokens": 1024
    }
    try:
        resp = requests.post(url, headers=headers, json=payload)
        resp.raise_for_status()
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception as exc:
        return f"[Error calling GPT API: {exc}]"

@app.post("/query_documents_with_page_range", response_model=QueryResponse)
def query_documents(req: QueryRequest):
    """
    1. Embeds the user prompt with mpnet
    2. Searches in FAISS
    3. Filters by the selected file(s) and page range(s)
    4. Takes top_k results
    5. Calls the GPT API with those top K as context
    """
    if faiss_index.ntotal == 0:
        return QueryResponse(answer="No data available. FAISS index is empty.", top_k_metadata=[])

    # 1) embed
    user_emb = generate_embeddings(req.prompt).reshape(1, -1)
    k = faiss_index.ntotal
    D, I = faiss_index.search(user_emb, k)

    # 2) filter
    relevant = []
    for dist, idx in zip(D[0], I[0]):
        if idx < len(metadata_store):
            record = metadata_store[idx]
            if record["filename"] in req.selected_files:
                (start_pg, end_pg) = req.selected_page_ranges.get(record["filename"], (1,999999))
                if start_pg <= record["page"] <= end_pg:
                    relevant.append((dist, record))

    # sort ascending by distance, take top_k
    relevant_sorted = sorted(relevant, key=lambda x: x[0])[: req.top_k]
    top_k_metadata = [r[1] for r in relevant_sorted]

    # 3) build final user prompt with the context
    context_json = json.dumps(top_k_metadata, indent=2)
    final_user_prompt = f"""
    The user query is: {req.prompt}

    Here are the top {req.top_k} matching segments from the user's documents:
    {context_json}

    Please use only this context to answer. 
    If something is not mentioned, say 'information not found.'
    """

    # 4) call your GPT function
    answer = call_gpt_api(SYSTEM_MESSAGE, final_user_prompt)

    return QueryResponse(answer=answer, top_k_metadata=top_k_metadata)
